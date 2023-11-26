import os
from typing import Union
from split_text import split_text_by_chars
from embedding import generate_embed_openai, language_embed
from translation import NLPTranslation
import tabula
import PyPDF2
import openai
import pandas as pd
import asyncio
import fitz
import base64
from docx import Document
from striprtf.striprtf import rtf_to_text
import ezodf
from pptx import Presentation
from odf import text, teletype
from odf.opendocument import load
import ebooklib
from ebooklib import epub
from bs4 import BeautifulSoup
import json
from openai_components import num_tokens_from_string, generate_chat_response

#######################################################################################################################################
########################################################## MAIN #######################################################################
#######################################################################################################################################

async def Embed_document(file_name, NLP_api_key, openai_api_key, language):
    # dictionary mapping file extensions to their respective functions
    file_extensions_to_functions = {
        '.pdf': embed_pdf,
        '.docx': embed_docx,
        '.txt': embed_txt,
        '.rtf': embed_rtf,
        '.odt': embed_odt,
        '.md': embed_md,
        '.pptx': embed_pptx,
        '.odp': embed_odp,
        '.epub': embed_epub,
    }

    _, file_extension = os.path.splitext(file_name)
    
    # Check if the file extension exists in the dictionary
    if file_extension in file_extensions_to_functions:
        # If it does, call the respective function
        output = await file_extensions_to_functions[file_extension](file_name, openai_api_key, NLP_api_key, language)
        return output
    else:
        # If it doesn't, raise an error
        # raise ValueError(f"File type {file_extension} is not supported.")
        return None

#######################################################################################################################################
################################################### For dictionary Structure ##########################################################
#######################################################################################################################################

async def create_text_content(page, text, english_text, embed):
    return {
        'class': "text",
        'page': page,
        'text': text,
        'english text': english_text,
        'embed': embed
    }

async def create_table_content(page, table, text, english_text, embed):
    return {
        'class': "table",
        'page': page,
        'table': table,
        'text': text,
        'english text': english_text,
        'embed': embed
    }

async def create_image_content(page, image, file_extension, text, english_text, embed):
    return {
        'class': "image",
        'page': page,
        'image': image,
        'file extension': file_extension,
        'text': text,
        'english text': english_text,
        'embed': embed
    }

async def create_result(file_name, file_extension, description, language, embed_language, content):
    return {
        'class': "document",
        'description': description,
        'file name': file_name,
        'file extension': file_extension,
        'language': {'prefix': language, 'ai language': '', 'embed language': embed_language},
        'content': content
    }


#######################################################################################################################################
########################################################### For text content ##########################################################
#######################################################################################################################################

async def extract_page_content(page_num, text, NLP_api_key, openai_api_key, language, embed_translate):
    texts_list = await split_text_by_chars(text, 1000)
    texts = texts_list
    
    # If translation is needed, translate the text into English.
    if embed_translate:
        try:
            english_texts = [await NLPTranslation(text_piece, NLP_api_key, language, "en") for text_piece in texts_list]
            texts = [response["translated text"] for response in english_texts]
        except Exception as e:
            raise RuntimeError(f"An error occurred while translating the text: {str(e)}")
        
    # Generate content for each text piece.
    page_content = []
    for i, text in enumerate(texts):
        embed_result = await generate_embed_openai(text, openai_api_key, language)

        english_text = text if embed_translate else ""
        original_text = texts_list[i] if embed_translate else text

        content_piece = await create_text_content(page_num, original_text, english_text, embed_result["embed"])
        page_content.append(content_piece)
    
    return page_content

#######################################################################################################################################
########################################################### For table content #########################################################
#######################################################################################################################################

async def extract_table_content(sheet, table, NLP_api_key, openai_api_key, language, embed_translate):
    # Validate if table meets the required conditions
    if not await is_dataframe_valid(table):
        return None
    
    # Convert the table (DataFrame) to a list of lists, including column headers.
    # This list format will be used for the final output.
    original_table = [table.columns.tolist()] + table.values.tolist()

    # Translate the table content if the 'embed_translate' flag is True
    translated_table = await translate_table(table, language, NLP_api_key) if embed_translate else table

    # Split the translated table into chunks if it exceeds a certain token limit
    chunks = await split_table_into_chunks(translated_table, token_limit=4096, column_multiplier=30)

    # Check if the full table is being used. If there is only one chunk, it implies that the full table is being used.
    full_table = len(chunks) == 1

    # Interpret each chunk of data and generate a summary and row descriptions
    output_split = await interpret_table_chunks(chunks, full_table)

    # Initialize text and english_text lists to store the processed output
    text, english_text, embeddings = [], [], []

    # Process each element in output_split
    for element in output_split:
        # Generate embeddings for the English text
        embed = await generate_embed_openai(element, openai_api_key, language) if element else {'embed': []}
        
        embeddings.append(embed["embed"])
        
        # If translation is needed, translate the text back to the original language and add to the lists
        if embed_translate:
            translated_text = (await NLPTranslation(element, NLP_api_key, "en", language))["translated text"]
            text.append(translated_text)
            english_text.append(element)
        else:
            text.append(element)
            english_text.append("")

    # Call create_table_content function to structure the final output
    table_content = await create_table_content(sheet, original_table, text, english_text, embeddings)
    return table_content

# Function to translate a table using an NLP API, if 'embed_translate' flag is True
async def translate_table(table, embed_translate, language, NLP_api_key):
    result = table.copy()
    
    # Generate translation tasks for non-null, non-empty and non-numeric cells
    tasks = [(row, col, NLPTranslation(cell, NLP_api_key, language, "en")) 
             for row in table.iterrows() for col in row[1].index 
             if pd.notnull(cell := row[1][col]) and cell != "" and not isinstance(cell, (int, float))]

    # Execute translation tasks and update the result table with translated content
    for row, col, task in tasks:
        result.at[row, col] = await task

    return result

# Function to validate if a DataFrame is valid based on certain conditions
async def is_dataframe_valid(df: pd.DataFrame, dot_threshold: float = 0.2, nan_threshold: float = 0.4) -> bool:
    # Checks whether the input DataFrame is valid based on two conditions:
    # 1. The percentage of periods ('.') in the string representation of the DataFrame should be below `dot_threshold`.
    # 2. The percentage of NaN values in the DataFrame should be below `nan_threshold`.

    # Concatenate all the data in the DataFrame into one string and remove spaces
    data_string_no_spaces = df.to_string().replace(" ", "")
    
    # Get the total number of cells in the DataFrame
    total_cells = df.size

    # If DataFrame is empty, return False
    if total_cells == 0 or len(data_string_no_spaces) == 0:
        return False

    # Calculate the percentage of periods and NaN values in the DataFrame
    period_percentage = data_string_no_spaces.count('.') / len(data_string_no_spaces)
    nan_percentage = df.isnull().sum().sum() / total_cells

    # Return True if both percentages are below their respective thresholds, otherwise False
    return period_percentage < dot_threshold and nan_percentage < nan_threshold

# Function to split the table into chunks if it exceeds a certain token limit
async def split_table_into_chunks(table, token_limit, column_multiplier):
    # This function works by calculating the number of tokens in each row of the table.
    
    # Create a string representing all the column names
    columns_str = "  ".join(str(column) for column in table.columns)

    # Initialize list to store token count of each row
    row_token_counts = [await num_tokens_from_string(columns_str, "cl100k_base")]

    # Calculate token count for each row and add to the list
    for _, row in table.iterrows():
        tokens = await num_tokens_from_string(str(row.values)[1:-1], "cl100k_base")
        row_token_counts.append(tokens)

    # Check if token count exceeds limit
    total_tokens = sum(row_token_counts) * (1 + column_multiplier / (table.shape[0] + table.shape[1]))
    high_token = total_tokens > token_limit

    # If it does, split the table into chunks
    if high_token:
        chunks = await create_table_chunks(table, row_token_counts, token_limit, column_multiplier)
    else:
        chunks = [table]
    return chunks

# Helper function to create chunks of table based on token limit
async def create_table_chunks(table, row_token_counts, token_limit, column_multiplier):
    # If the total number of tokens exceeds the token limit, it splits the table into chunks.

    chunks, chunk, current_token_count = [], [], 0
    for idx, row in enumerate(table.itertuples(index=False)):
        row_token = row_token_counts[idx + 1] * (1 + column_multiplier / table.shape[1])
        if current_token_count + row_token > token_limit and chunk:
            chunks.append(pd.DataFrame(chunk, columns=table.columns))
            chunk, current_token_count = [], 0
        chunk.append(dict(zip(table.columns, row)))
        current_token_count += row_token
    if chunk:
        chunks.append(pd.DataFrame(chunk, columns=table.columns))
    return chunks

# Function to interpret each chunk of the table, generate summary and row descriptions
async def interpret_table_chunks(chunks, full_table):
    # Ensure chunks is a list of pandas DataFrames
    if not isinstance(chunks, list) or not all(isinstance(chunk, pd.DataFrame) for chunk in chunks):
        raise TypeError("chunks must be a list of pandas DataFrame objects.")

    # Ensure full_table is a boolean
    if not isinstance(full_table, bool):
        raise TypeError("full_table must be a boolean.")

    # Initialize output list
    more_output = []

    # Interpret each chunk
    for chunk_df in chunks:
        chunk_string = await generate_chunk_string(chunk_df)
        
        messages = [
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": f"""Use the table provided below to analyze the data and gain valuable insights. Focus on the numbers and identify trends that reveal important information. Highlight the most significant rows to draw attention to key findings. Create a narrative that is engaging and optimized for semantic search, ensuring that users can easily find the information they need.

Instead of using lists or bullet points, craft well-connected paragraphs that flow together seamlessly. Paint a clear and vivid picture of the data, captivating readers and encouraging them to explore further. Consider the possible questions users might ask and incorporate them naturally into your text to enhance its visibility in semantic search results.

Table: {chunk_string}"""},
        ]

        try:
            output = await generate_chat_response(messages)
            output = output["content"]
            output_split = await split_text_by_chars(output, 1000)
            
            if full_table:
                return output_split
            more_output.extend(output_split)
        except Exception as e:
            raise Exception(f"An error occurred: {e}")
    return more_output

# Helper function to generate string representation of a table chunk
async def generate_chunk_string(chunk_df):
    chunk_str = f"Column's name': {', '.join(str(chunk_df.columns))}\n"
    for i, row in chunk_df.iterrows():
        row_str = [str(element) for element in row]
        chunk_str += f"Row {i+1}: {', '.join(row_str)}\n"
    return chunk_str

#######################################################################################################################################
########################################################### For image content #########################################################
#######################################################################################################################################

async def extract_image_content(page, image_data, image_extension, NLP_api_key, openai_api_key, language, embed_translate):
    # Initialize the English text with a placeholder value, to be replaced with the actual image description.
    english_text = "nothing"

    # Translate the English text to the specified language using the NLPTranslation function.
    # translated_text_response = await NLPTranslation(english_text, NLP_api_key, 'en', language)
    # translated_text = translated_text_response["translated text"]

    # Decide which text to use for embedding based on the embed_translate flag.
    # text_to_embed = english_text if embed_translate else translated_text

    # Generate the embedded result using the OpenAI API.
    # embed_result = await generate_embed_openai(text_to_embed, openai_api_key, language)
    
    # If the specified language is English, clear the translated text variable.
    # if language == 'en':
    #    english_text, translated_text = translated_text, ""
    
    translated_text = ""
    embed_result = []

    # Create the image content using the create_image_content function.
    content = await create_image_content(page, image_data, image_extension, translated_text, english_text, embed_result["embed"])

    return content

#######################################################################################################################################
################################################################## PDF ################################################################
#######################################################################################################################################

async def embed_pdf(file_path, openai_api_key, NLP_api_key, language):
    try:
        openai.api_key = openai_api_key

        pdf = fitz.open(file_path)
        pdf_reader = PyPDF2.PdfReader(open(file_path, 'rb'))

        embed_translate = await language_embed(language)
        embed_language = 'en' if embed_translate else language

        content = {}
        content_id = 0  

        for page_num, page in enumerate(pdf_reader.pages):
            raw_text = page.extract_text()
            print(raw_text)
            print()
            page_content = await extract_page_content(page_num, raw_text, NLP_api_key, openai_api_key, language, embed_translate)

            for content_piece in page_content:
                content[content_id] = content_piece
                content_id += 1

            page_tables = tabula.read_pdf(file_path, pages=page_num + 1)
            for table in page_tables:
                interpretation = await extract_table_content(page_num, table, NLP_api_key, openai_api_key, language, embed_translate)
                if interpretation is not None:
                    content[content_id] = interpretation
                    content_id += 1

        for page_num, page in enumerate(pdf):
            images = page.get_images()
            images.reverse()
        
            for image in images:
                base_image = pdf.extract_image(image[0])
                image_data = base_image["image"]
                image_string = base64.b64encode(image_data).decode('utf-8')
                image_extension = base_image["ext"]
                interpretation = await extract_image_content(page_num, image_string, image_extension, NLP_api_key, openai_api_key, language, embed_translate)
                if interpretation is not None:
                    content[content_id] = interpretation
                    content_id += 1

        file_name = os.path.basename(file_path)
        file_extension = os.path.splitext(file_name)[1]
        description = ""
        document = await create_result(file_name, file_extension, description, language, embed_language, content)
        
        return document
    except Exception as e:
        raise Exception(f"An error occurred: {e}")

#######################################################################################################################################
################################################################## DOCX ###############################################################
#######################################################################################################################################

async def embed_docx(file_path, openai_api_key, NLP_api_key, language = "en") -> Union[dict, None]:
    openai.api_key = openai_api_key

    try:
        # Open the document
        doc = Document(file_path)
    except Exception as e:
        raise FileNotFoundError("Could not open file. Please check the file path.")

    # Decide which language to use for embedding
    embed_translate = await language_embed(language)
    embed_language = 'en' if embed_translate else language

    content = {}  
    content_id = 0  # Initialize counter for unique content IDs
    raw_text = ""
    page_num = -1  # No page structure

    # Extract and process text content
    for para in doc.paragraphs:
        raw_text += para.text
    
    text_contents = await extract_page_content(page_num, raw_text, NLP_api_key, openai_api_key, language, embed_translate)

    for text_content in text_contents:
        content[content_id] = text_content
        content_id += 1  # Increment content ID
        
    # Extract and process table content
    for table in doc.tables:
        data = []
        keys = None
        for i, row in enumerate(table.rows):
            text = (cell.text for cell in row.cells)
            if i == 0:  # If it's the first row, consider it as keys
                keys = list(text)
            else:
                row_data = dict(zip(keys, text))  # Create a dictionary for each row
                data.append(row_data)
            
        df = pd.DataFrame(data)
        
        page_num = 0  # Initialize page number for tables
        interpretation = await extract_table_content(page_num, df, NLP_api_key, openai_api_key, language, embed_translate)
        if interpretation is not None:
            content[content_id] = interpretation
            content_id += 1  # Increment content ID
    
    # Extract and process image content
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image_data = rel._target._blob
            image_string = base64.b64encode(image_data).decode('utf-8')  # Convert image data to base64
            image_extension = rel.reltype.split('.')[-1]  # Extract image extension from relationship type
            interpretation = await extract_image_content(page_num, image_string, image_extension, NLP_api_key, openai_api_key, language, embed_translate)
            if interpretation is not None:
                content[content_id] = interpretation
                content_id += 1  # Increment content ID

    # Get file name and extension
    file = os.path.basename(file_path)
    file_extension = os.path.splitext(file)[1]
    description = ""
    document = await create_result(file, file_extension, description, language, embed_language, content)
    
    return document

#######################################################################################################################################
################################################################### RTF ###############################################################
#######################################################################################################################################

async def embed_rtf(file_path, openai_api_key, NLP_api_key, language = "en"):
    with open(file_path, 'r') as file:
        rtf_text = file.read()
    plain_text = rtf_to_text(rtf_text)
        
    combined_text, list_df = text_to_table_text(plain_text)
    
    content = {}
    idx = 0
    
    embed_translate = await language_embed(language)
    if embed_translate:
        embed_language = 'en'
    else: 
        embed_language = language
    content = {}
    content_id = 0  # Initialize counter for unique content IDs
    
    page_num = -1

    # Extract and process text content
    raw_text = combined_text
    page_content = await extract_page_content(page_num, raw_text, NLP_api_key, openai_api_key, language, embed_translate)
    
    for content_piece in page_content:
        content[content_id] = content_piece
        content_id += 1
    
    for table in list_df:
        page_num = 0
        interpretation = await extract_table_content(page_num, table, NLP_api_key, openai_api_key, language, embed_translate)
        content[idx] = interpretation
        
    file = os.path.basename(file_path)
    file_extension = os.path.splitext(file)[1]
    description = ""
    document = await create_result(file, file_extension, description, language, embed_language, content)
    
    return document

def text_to_table_text(table_str) -> pd.DataFrame:
    lines = table_str.strip().split("\n")

    tables = []
    now = []
    texts = []
    for line in lines:
        if "|" in line:
            now.append(line)
        else:
            if len(now) > 1:
                tables.append(now)
                now = []
            texts.append(line)
            
    list_df = []
    for data in tables:
        df = pd.DataFrame([x.split('|') for x in data], columns=data[0].split('|'))
        df = df.iloc[1:]  # Remove the first row as it contains column names
        list_df.append(df)
    
    combined_text = " ".join(texts).replace("  ", "\n")

    return combined_text, list_df

#######################################################################################################################################
################################################################### ODT ###############################################################
#######################################################################################################################################

async def embed_odt(file_path, openai_api_key, NLP_api_key, language):
    doc = ezodf.opendoc(file_path)

    # Extract text
    text = []
    for paragraph in doc.body:
        if isinstance(paragraph, ezodf.text.Paragraph):
            text.append(paragraph.plaintext())

    combined_text = ' '.join(text)
    
    content = {}
    idx = 0
    page_num = -1
    
    embed_translate = await language_embed(language)
    if embed_translate:
        embed_language = 'en'
    else: 
        embed_language = language
    content = {}
    content_id = 0  # Initialize counter for unique content IDs

    # Extract and process text content
    page_content = await extract_page_content(page_num, combined_text, NLP_api_key, openai_api_key, language, embed_translate)
    
    for content_piece in page_content:
        content[content_id] = content_piece
        content_id += 1
    
    # table
    list_df = []
            
    for table in doc.body:
        if isinstance(table, ezodf.table.Table):
            tables = []
            for row in table:
                row_data = []
                for cell in row:
                    row_data.append(cell.plaintext())
                if row_data != []:
                    tables.append(row_data)
            
            df = pd.DataFrame(tables)
            df.columns = df.iloc[0]  # Set the first (0th) row as header
            df = df.iloc[1:]  # Remove the first row
            list_df.append(df)
    
    for table in list_df:
        page_num = 0
        interpretation = await extract_table_content(page_num, table, NLP_api_key, openai_api_key, language, embed_translate)
        content[idx] = interpretation
        
    file = os.path.basename(file_path)
    file_extension = os.path.splitext(file)[1]
    description = ""
    document = await create_result(file, file_extension, description, language, embed_language, content)
    
    return document

#######################################################################################################################################
################################################################### TXT ###############################################################
#######################################################################################################################################

async def embed_txt(file_path, openai_api_key, NLP_api_key, language):
    with open(file_path, 'r', encoding='utf-8') as file:
        text = file.read()

    embed_translate = await language_embed(language)
    if embed_translate:
        embed_language = 'en'
    else: 
        embed_language = language
    content = {}
    content_id = 0  # Initialize counter for unique content IDs

    page_num = -1
    page_content = await extract_page_content(page_num, text, NLP_api_key, openai_api_key, language, embed_translate)

    for content_piece in page_content:
        content[content_id] = content_piece
        content_id += 1

    file_name = os.path.basename(file_path)
    file_extension = os.path.splitext(file_name)[1]
    description = ""
    document = await create_result(file_name, file_extension, description, language, embed_language, content)
    
    return document

#######################################################################################################################################
################################################################### MD ################################################################
#######################################################################################################################################

async def embed_md(file_path, openai_api_key, NLP_api_key, language):
    with open(file_path, 'r', encoding='utf-8') as file:
        text = file.read()

    embed_translate = await language_embed(language)
    if embed_translate:
        embed_language = 'en'
    else: 
        embed_language = language
        
    content = {}
    content_id = 0  # Initialize counter for unique content IDs

    page_num = -1
    page_content = await extract_page_content(page_num, text, NLP_api_key, openai_api_key, language, embed_translate)

    for content_piece in page_content:
        content[content_id] = content_piece
        content_id += 1

    file_name = os.path.basename(file_path)
    file_extension = os.path.splitext(file_name)[1]
    description = ""
    document = await create_result(file_name, file_extension, description, language, embed_language, content)
    
    return document
    
#######################################################################################################################################
############################################################# PPTX ####################################################################
#######################################################################################################################################

async def embed_pptx(file_path, openai_api_key, NLP_api_key, language):
    openai.api_key = openai_api_key

    # Open the pptx file.
    presentation = Presentation(file_path)
    
    embed_translate = await language_embed(language)
    if embed_translate:
        embed_language = 'en'
    else: 
        embed_language = language
        
    content = {}
    content_id = 0  # Initialize counter for unique content IDs
    
    for page_num, slide in enumerate(presentation.slides):
        # Extract and process text content
        raw_text = ' '.join([shape.text for shape in slide.shapes if shape.has_text_frame])
        page_content = await extract_page_content(page_num, raw_text, NLP_api_key, openai_api_key, language, embed_translate)
        
        for content_piece in page_content:
            content[content_id] = content_piece
            content_id += 1

        # Extract and process table content
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                table_data = [[cell.text for cell in row.cells] for row in table.rows]
                
                # Remove empty rows from the top of the table
                while table_data and all(cell == '' for cell in table_data[0]):
                    table_data = table_data[1:]
        
                if table_data:  # Skip tables that have no rows left after the removal
                    df = pd.DataFrame(table_data)
                    df.columns = df.iloc[0]  # Make the first row as header
                    df = df[1:]  # Remove the first row
                    interpretation = await extract_table_content(page_num, df, NLP_api_key, openai_api_key, language, embed_translate)
                    if interpretation is not None:
                        content[content_id] = interpretation
                        content_id += 1

        # Extract and process image content
        for shape in slide.shapes:
            if shape.shape_type == 13:  # This is the code for a Picture
                image = shape.image
                image_bytes = image.blob
                image_string = base64.b64encode(image_bytes).decode('utf-8')
                interpretation = await extract_image_content(page_num, image_string, 'png', NLP_api_key, openai_api_key, language, embed_translate)
                if interpretation is not None:
                    content[content_id] = interpretation
                    content_id += 1

    file = os.path.basename(file_path)
    file_extension = os.path.splitext(file)[1]
    description = ""
    document = await create_result(file, file_extension, description, language, embed_language, content)
    
    return document

#######################################################################################################################################
############################################################## ODP ####################################################################
#######################################################################################################################################

async def embed_odp(file_path, openai_api_key, NLP_api_key, language):
    openai.api_key = openai_api_key

    # Load the ODP file.
    presentation = load(file_path)
    
    embed_translate = await language_embed(language)
    if embed_translate:
        embed_language = 'en'
    else: 
        embed_language = language
        
    content = {}
    content_id = 0  # Initialize counter for unique content IDs    
    
    # No slide structure
    page_num = -1
    
    raw_text = ""

    for slide in presentation.getElementsByType(text.P):
        # Extract and process text content
        raw_text = raw_text + "\n\n" + teletype.extractText(slide)

    page_content = await extract_page_content(page_num, raw_text, NLP_api_key, openai_api_key, language, embed_translate)
        
    for content_piece in page_content:
        content[content_id] = content_piece
        content_id += 1

    file = os.path.basename(file_path)
    file_extension = os.path.splitext(file)[1]
    description = ""
    document = await create_result(file, file_extension, description, language, embed_language, content)
    
    return document

#######################################################################################################################################
############################################################# EPUB ####################################################################
#######################################################################################################################################

async def embed_epub(file_path, openai_api_key, NLP_api_key, language):
    openai.api_key = openai_api_key

    # Open the epub file.
    book = epub.read_epub(file_path)
    items = list(book.get_items_of_type(ebooklib.ITEM_DOCUMENT))

    embed_translate = await language_embed(language)
    if embed_translate:
        embed_language = 'en'
    else: 
        embed_language = language

    content = {}
    content_id = 0  # Initialize counter for unique content IDs
    
    page_num = -1
    
    raw_text_combined = ""

    # Iterate through the items in the epub file.
    for item in items:
        # Extract and process text content
        soup = BeautifulSoup(item.get_body_content(), 'html.parser')
        raw_text = ' '.join(str(para.get_text()) for para in soup.find_all('p'))
        raw_text_combined = raw_text_combined + " " + raw_text
        
    page_content = await extract_page_content(page_num, raw_text_combined, NLP_api_key, openai_api_key, language, embed_translate)
    
    for content_piece in page_content:
        content[content_id] = content_piece
        content_id += 1

    # Epub doesn't contain tables like PDF, so no need to process tables.

    # Process images
    images = list(book.get_items_of_type(ebooklib.ITEM_IMAGE))
    for image in images:
        image_data = image.get_content()
        image_string = base64.b64encode(image_data).decode('utf-8')
        # Image extension can be determined from the item's file_name.
        image_extension = os.path.splitext(image.file_name)[1]
        interpretation = await extract_image_content(page_num, image_string, image_extension, NLP_api_key, openai_api_key, language, embed_translate)
        if interpretation is not None:
            content[content_id] = interpretation
            content_id += 1


    file = os.path.basename(file_path)
    file_extension = os.path.splitext(file)[1]
    description = ""
    document = await create_result(file, file_extension, description, language, embed_language, content)
    
    return document

#######################################################################################################################################
########################################################### RUN TESTING AREA ##########################################################
#######################################################################################################################################


# To call the function
async def call():
    openai_api_key = "openai_api_key"
    NLP_api_key = "NLP_api_key"
    file_name = "fileNameToBreakDown"
    language = "en"
    output_dict = await Embed_document(file_name, openai_api_key, NLP_api_key, language)
        
    # Define the output file path
    output_file = "output.json"
    
    # Write the dictionary to a JSON file
    with open(output_file, "w") as f:
        json.dump(output_dict, f)
    
asyncio.run(call())















